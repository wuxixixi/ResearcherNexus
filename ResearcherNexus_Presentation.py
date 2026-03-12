from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 重命名RGBColor为RgbColor以保持代码一致性
RgbColor = RGBColor

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色方案
PRIMARY_COLOR = RgbColor(0x1E, 0x3A, 0x5F)      # 深蓝
ACCENT_COLOR = RgbColor(0x3B, 0x82, 0xF6)       # 亮蓝
TEXT_COLOR = RgbColor(0x1F, 0x29, 0x37)         # 深灰
LIGHT_BG = RgbColor(0xF3, 0xF4, 0xF6)           # 浅灰背景
WHITE = RgbColor(0xFF, 0xFF, 0xFF)

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 添加背景形状
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()

    # 装饰线条
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.5), Inches(2), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(0xBF, 0xDB, 0xFE)

    return slide

def add_section_slide(prs, title, subtitle=""):
    """添加章节分隔页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(11), Inches(0.6))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = RgbColor(0xBF, 0xDB, 0xFE)

    return slide

def add_content_slide(prs, title, bullets, icon=""):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 浅色背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.fill.background()

    # 左侧装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # 内容区域
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12), Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # 支持子项目（以 "  - " 开头）
        if bullet.startswith("  - "):
            p.text = "• " + bullet[4:]
            p.level = 1
            p.font.size = Pt(16)
            p.font.color.rgb = RgbColor(0x4B, 0x55, 0x63)
        else:
            p.text = "● " + bullet
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(10)

    return slide

def add_two_column_slide(prs, title, left_title, left_content, right_title, right_content):
    """添加双栏内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.fill.background()

    # 装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # 左栏
    left_title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(5.8), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    left_content_box = slide.shapes.add_textbox(Inches(0.6), Inches(2), Inches(5.8), Inches(5))
    tf = left_content_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "● " + line
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(6)

    # 右栏
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(5.8), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    right_content_box = slide.shapes.add_textbox(Inches(6.8), Inches(2), Inches(5.8), Inches(5))
    tf = right_content_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "● " + line
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(6)

    return slide

def add_architecture_slide(prs):
    """添加架构图页面"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.fill.background()

    # 装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "系统技术架构"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # 架构层次 - 使用表格方式展示
    y_start = 1.4
    box_height = 0.9
    box_width = 12

    layers = [
        ("Web 前端 (Next.js + React)", "用户界面、可视化编辑、流程图展示", RgbColor(0xDB, 0xEA, 0xFE)),
        ("后端 API (FastAPI)", "RESTful接口、WebSocket通信、会话管理", RgbColor(0xE0, 0xE7, 0xFF)),
        ("LangGraph 工作流引擎", "多代理系统、状态管理、任务调度", RgbColor(0xFE, 0xE2, 0xE2)),
        ("工具集成层", "搜索引擎、爬虫、代码执行、MCP工具", RgbColor(0xF3, 0xE8, 0xFF)),
        ("语言模型层", "OpenAI API、开源模型、LLM统一接口", RgbColor(0xFF, 0xF7, 0xED)),
    ]

    for i, (layer_name, desc, color) in enumerate(layers):
        y = y_start + i * 1.0

        # 主框
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.6), Inches(y),
                                     Inches(box_width), Inches(box_height))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = RgbColor(0x9C, 0xA3, 0xAF)
        box.line.width = Pt(1)

        # 层名称
        name_box = slide.shapes.add_textbox(Inches(0.9), Inches(y + 0.15), Inches(3), Inches(0.5))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = layer_name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEXT_COLOR

        # 描述
        desc_box = slide.shapes.add_textbox(Inches(4.2), Inches(y + 0.2), Inches(7.8), Inches(0.5))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(0x4B, 0x55, 0x63)

    return slide

# ==================== 开始创建PPT ====================

# 1. 封面页
add_title_slide(prs, "ResearcherNexus", "AI驱动的智能研究自动化平台")

# 2. 目录页
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.line.fill.background()

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_COLOR
bar.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "目录"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = PRIMARY_COLOR

toc_items = [
    "01  项目概述与核心价值",
    "02  系统架构与技术栈",
    "03  核心功能模块",
    "04  产品特色与优势",
    "05  应用场景与价值",
    "06  项目成果与展望"
]

for i, item in enumerate(toc_items):
    y = 1.5 + i * 0.9
    item_box = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(12), Inches(0.6))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT_COLOR

# 3. 项目概述 - 章节页
add_section_slide(prs, "01  项目概述与核心价值", "PROJECT OVERVIEW")

# 4. 什么是ResearcherNexus
add_content_slide(prs, "什么是 ResearcherNexus？", [
    "ResearcherNexus 是一个社区驱动的深度研究自动化框架",
    "基于多代理协作系统，将语言模型与专业研究工具相结合",
    "支持网络搜索、网页爬取、Python代码执行等多种研究手段",
    "采用 LangGraph 工作流引擎，实现基于状态的智能任务调度",
    "支持生成研究报告、播客、演示文稿等多模态内容输出",
    "具备人机协作能力，支持自然语言交互修改研究计划"
])

# 5. 核心价值主张
add_two_column_slide(prs, "核心价值主张",
    "为研究者赋能",
    [
        "自动化繁琐的研究资料收集",
        "快速整合多源异构信息",
        "智能生成结构化研究报告",
        "支持多格式内容输出",
        "大幅降低研究时间成本"
    ],
    "为企业创造价值",
    [
        "提升市场分析效率",
        "加速技术趋势研判",
        "支持投资决策研究",
        "标准化研究流程",
        "积累企业知识资产"
    ]
)

# 6. 系统架构 - 章节页
add_section_slide(prs, "02  系统架构与技术栈", "SYSTEM ARCHITECTURE")

# 7. 系统架构图
add_architecture_slide(prs)

# 8. 后端技术栈
add_content_slide(prs, "后端技术栈", [
    "Python 3.12+ - 核心编程语言",
    "LangGraph 0.3.5 - 工作流引擎与状态管理",
    "LangChain 0.3.x - LLM集成与工具链",
    "FastAPI 0.110.0 - 高性能API框架",
    "LiteLLM 1.63.11 - 多LLM统一接口",
    "MCP 1.6.0 - Model Context Protocol支持",
    "Tavily/DuckDuckGo/Brave - 多搜索引擎集成"
])

# 9. 前端技术栈
add_content_slide(prs, "前端技术栈", [
    "Next.js 15.3.2 - React全栈框架",
    "React 19.0.0 - 用户界面库",
    "TypeScript 5.8.2 - 类型安全开发",
    "Tailwind CSS 4.0.15 - 原子化样式框架",
    "Tiptap/ProseMirror - 富文本编辑器",
    "xyflow - 工作流可视化引擎",
    "Zustand - 轻量级状态管理"
])

# 10. 核心功能模块 - 章节页
add_section_slide(prs, "03  核心功能模块", "CORE MODULES")

# 11. 多代理协作系统
add_content_slide(prs, "多代理协作系统 (Multi-Agent System)", [
    "协调器 (Coordinator)：管理工作流生命周期，启动和监控研究过程",
    "规划器 (Planner)：分析研究目标，创建结构化执行计划",
    "研究员 (Researcher)：执行网络搜索、信息收集和资料整理",
    "编码员 (Coder)：处理代码分析、执行和Python数据处理任务",
    "报告员 (Reporter)：汇总研究发现，生成结构化研究报告",
    "各代理通过LangGraph状态机协作，支持并行执行和条件分支"
])

# 12. 工具集成系统
add_two_column_slide(prs, "工具集成系统",
    "搜索工具",
    [
        "Tavily Search - AI专用搜索API",
        "DuckDuckGo - 隐私搜索引擎",
        "Brave Search - 高级搜索功能",
        "Arxiv - 学术论文搜索"
    ],
    "其他工具",
    [
        "Jina/ReadabiliPy - 网页内容提取",
        "Python REPL - 代码执行环境",
        "MCP协议 - 外部工具集成",
        "Volcengine TTS - 语音合成"
    ]
)

# 13. MCP协议集成
add_content_slide(prs, "MCP (Model Context Protocol) 集成", [
    "MCP是Anthropic推出的开放协议，用于标准化AI模型与外部工具的连接",
    "支持扩展私有域访问、知识图谱、网页浏览等高级能力",
    "促进多样化研究工具和方法的无缝集成",
    "内置智能工具选择和推荐算法，自动选择最适合的工具组合",
    "针对Windows环境进行兼容性优化，支持企业部署场景",
    "工具兼容性检查和错误恢复机制，确保系统稳定性"
])

# 14. 内容生成模块
add_content_slide(prs, "多模态内容生成", [
    "研究报告生成 - 自动生成结构化学术报告，支持多种引用格式",
    "播客生成 (Podcast) - 将研究报告转换为播客脚本并生成音频文件",
    "演示文稿生成 (PPT) - 自动创建PowerPoint演示文稿，支持图表和数据可视化",
    "散文/报告生成 (Prose) - 高级文本处理、润色和格式转换",
    "支持多种报告风格：学术、新闻、社交媒体、投资分析等",
    "引用管理系统 - 自动化引用收集、去重、格式化和报告标注"
])

# 15. 产品特色 - 章节页
add_section_slide(prs, "04  产品特色与优势", "KEY ADVANTAGES")

# 16. 产品特色
add_content_slide(prs, "产品核心特色", [
    "智能工作流系统 - 基于LangGraph的可视化工作流引擎，支持复杂条件分支和并行执行",
    "多代理协作 - 专门化代理角色，智能任务交接，支持自然语言指令和反馈",
    "MCP协议支持 - 无缝集成外部工具，支持私有域知识和企业工具，高度定制化",
    "高级引用管理 - 自动化引用收集和格式化，支持多种学术引用风格，符合学术标准",
    "多风格内容生成 - 支持多种报告风格，满足不同用户群体需求，提供专业级内容质量",
    "完善用户系统 - 完整用户认证授权，使用次数限制，管理员后台，详细使用统计"
])

# 17. 竞争优势
add_two_column_slide(prs, "竞争优势对比",
    "ResearcherNexus",
    [
        "多代理协作架构",
        "内置MCP协议支持",
        "可视化工作流引擎",
        "多模态内容生成",
        "开源+企业级功能",
        "灵活的部署选项"
    ],
    "传统研究工具",
    [
        "单一AI模型调用",
        "有限的工具集成",
        "线性处理流程",
        "纯文本输出为主",
        "闭源商业方案",
        "固定部署模式"
    ]
)

# 18. 应用场景 - 章节页
add_section_slide(prs, "05  应用场景与价值", "USE CASES")

# 19. 应用场景
add_content_slide(prs, "典型应用场景", [
    "学术研究与论文撰写 - 自动化文献综述、研究资料收集、引用管理、报告生成",
    "市场分析与竞争情报 - 行业趋势分析、竞品研究、市场机会识别、投资分析报告",
    "技术趋势与创新研究 - 前沿技术追踪、专利分析、技术路线图生成、创新机会评估",
    "新闻与内容创作 - 热点话题追踪、多角度信息整合、新闻稿生成、内容策划",
    "教育培训内容开发 - 课程资料准备、教学案例研究、培训材料生成、知识库建设",
    "企业知识管理 - 内部知识整合、专家经验提取、文档自动生成、知识图谱构建"
])

# 20. 业务价值
add_two_column_slide(prs, "业务价值",
    "效率提升",
    [
        "研究时间缩短70-90%",
        "信息收集自动化",
        "报告生成一键完成",
        "多任务并行处理",
        "7×24小时不间断工作"
    ],
    "质量保障",
    [
        "多源交叉验证",
        "自动引用标注",
        "结构化内容输出",
        "减少人为疏漏",
        "标准化研究流程"
    ]
)

# 21. 项目成果 - 章节页
add_section_slide(prs, "06  项目成果与展望", "ACHIEVEMENTS & ROADMAP")

# 22. 项目成果
add_content_slide(prs, "当前项目成果", [
    "完整的多代理协作系统 - 5种专业代理角色，基于LangGraph的工作流引擎",
    "丰富的工具集成 - 支持4种搜索引擎、2种网页爬取工具、Python代码执行环境",
    "MCP协议完整支持 - 标准化外部工具集成，智能工具推荐算法",
    "多模态内容生成 - 研究报告、播客、PPT、散文等多种输出格式",
    "完善的用户系统 - 认证授权、使用限制、管理员后台、使用统计",
    "现代化的前端界面 - Next.js + React，支持可视化编辑和流程图展示",
    "完整的文档和示例 - 配置指南、FAQ、9个完整示例报告"
])

# 23. 技术路线图
add_content_slide(prs, "技术发展路线图", [
    "近期目标 (v0.2.0) - 完整的引用管理系统、高级数据分析和可视化、更多LLM支持",
    "中期目标 (v0.3.0) - 高级协作功能、企业级部署选项、性能优化和扩展性提升",
    "长期愿景 (v1.0.0) - 自主研究能力、多模态理解、知识图谱集成、智能推荐系统",
    "持续优化方向 - 提升研究质量、扩展工具生态、增强用户体验、强化安全保障"
])

# 24. 结束页
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 背景
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = PRIMARY_COLOR
shape.line.fill.background()

# 装饰圆圈
for i, (x, y, size, opacity) in enumerate([(9, 1, 3, 0.1), (10, 4, 4, 0.08), (1, 5, 2, 0.1)]):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    circle.fill.solid()
    circle.fill.fore_color.rgb = WHITE
    circle.line.fill.background()

# 主标题
title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "ResearcherNexus"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = WHITE

# 副标题
sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11), Inches(0.8))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "让研究更智能 · 让知识更易得"
p.font.size = Pt(28)
p.font.color.rgb = RgbColor(0xBF, 0xDB, 0xFE)

# 联系信息
contact_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11), Inches(0.6))
tf = contact_box.text_frame
p = tf.paragraphs[0]
p.text = "社区驱动 · 开源开放 · 持续创新"
p.font.size = Pt(18)
p.font.color.rgb = RgbColor(0x9C, 0xA3, 0xAF)

# 保存PPT
output_path = r"D:\ResearcherNexus\ResearcherNexus_产品介绍.pptx"
prs.save(output_path)
print(f"PPT已保存至: {output_path}")
print(f"共 {len(prs.slides)} 页幻灯片")
